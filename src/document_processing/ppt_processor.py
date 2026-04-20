# -*- coding: utf-8 -*-
from __future__ import annotations

import os
from typing import Any, Dict, List

from .processor_utils import stable_doc_id


class PPTProcessor:
    """PPTX ingestion via python-pptx (PPT treated as unsupported)."""

    def process(self, abs_path: str, rel_path: str, export_root: str, *, doc_id_override: str | None = None) -> Dict[str, Any]:
        doc_id = doc_id_override or stable_doc_id(abs_path, prefix="pptx")
        ext = os.path.splitext(abs_path)[1].lower()

        pages: List[Dict[str, Any]] = []
        full_text_parts: List[str] = []

        if ext == ".pptx":
            try:
                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE

                prs = Presentation(abs_path)
                for i, slide in enumerate(prs.slides):
                    parts: List[str] = []
                    
                    if slide.shapes.title and slide.shapes.title.text:
                        parts.append(f"[Title] {slide.shapes.title.text.strip()}")

                    for shape in slide.shapes:
                        try:
                            if shape == slide.shapes.title:
                                continue
                                
                            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                                parts.append("[Table]")
                                for row in shape.table.rows:
                                    row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                                    parts.append(" | ".join(row_data))
                            elif hasattr(shape, "text") and shape.text:
                                t = str(shape.text).strip()
                                if t:
                                    parts.append(f"[Body] {t}")
                            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                try:
                                    image_bytes = shape.image.blob
                                    img_path = os.path.join(export_root, f"slide_{i}_img_{shape.shape_id}.png")
                                    with open(img_path, "wb") as f:
                                        f.write(image_bytes)
                                    from src.vision.ocr_backends import TesseractBackend
                                    ocr_text = TesseractBackend().recognize(img_path, lang_hint="eng")
                                    if ocr_text:
                                        parts.append(f"[Image OCR] {ocr_text.strip()}")
                                except Exception:
                                    pass
                        except Exception:
                            continue
                            
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            parts.append(f"[Speaker Notes] {notes}")

                    slide_text = "\n".join(parts).strip()
                    if slide_text:
                        full_text_parts.append(slide_text)

                    pages.append(
                        {
                            "page_index": i,
                            "py_text": slide_text,
                            "text_hint": (slide_text[:200] if slide_text else f"Slide {i+1}"),
                            "quality": "queued",
                        }
                    )
            except Exception as e:
                pages = [
                    {
                        "page_index": 0,
                        "py_text": f"[pptx] failed to parse: {e}",
                        "text_hint": rel_path,
                        "quality": "queued",
                    }
                ]
        else:
            pages = [
                {
                    "page_index": 0,
                    "py_text": "[ppt] legacy .ppt detected; conversion not implemented in this build.",
                    "text_hint": rel_path,
                    "quality": "queued",
                }
            ]

        return {
            "doc_id": doc_id,
            "text": "\n\n".join(full_text_parts).strip(),
            "pages": pages if pages else [{"page_index": 0, "py_text": "", "quality": "queued"}],
            "meta": {"source": "ppt-processor", "rel_path": rel_path, "slides": len(pages)},
        }

